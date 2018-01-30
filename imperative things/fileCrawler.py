#Amanda Doyle - 212464575#

from pptx import Presentation
import os
import csv
import re
from cStringIO import StringIO
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams
from pdfminer.pdfpage import PDFPage
import sys, getopt
from docx import Document
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML
import zipfile
import unicodedata

WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'

if os.path.exists("C:\Python27\crawlerOutput.txt"):
    os.remove("C:\Python27\crawlerOutput.txt")

text_file = open("C:\Python27\crawlerOutput.txt", "w")

files = os.listdir("C:\Python27\\files_to_parse\\")
ppt_files = []
pdf_files = []
doc_files = []
temp = []
words = []
counter = []
substring_counter = 0
number_of_docs = 0
file_names_to_read = []

def all_same(items):
    return all(x == items[0] for x in items)

for items in words:
	counter.append(0)

with open('keywords.csv', 'rb') as f:
    readCSV = csv.reader(f, delimiter=',')
    for row in readCSV:
    	words.append(row[0])

for names in files:
    if names.endswith(".ppt") or names.endswith(".pptx"):
        ppt_files.append(names)
    elif names.endswith(".docx"):
    	doc_files.append(names)
    elif names.endswith(".pdf"):
    	pdf_files.append(names)

print("\nParsing "+str(len(ppt_files))+ " PPT/PPTX files.")
print("Parsing "+str(len(pdf_files))+ " PDF files.")
print("Parsing "+str(len(doc_files))+ " DOCX files.\n")

text_file.write("\nParsed "+str(len(ppt_files))+ " PPT/PPTX files.\n")
text_file.write("Parsed "+str(len(pdf_files))+ " PDF files.\n")
text_file.write("Parsed "+str(len(doc_files))+ " DOCX files.\n\n")

array = []

#DOCX Parser#

def getDocxText(filename):
	array[:] = []
	document = Document(filename)
	for p in document.paragraphs:
		if p.text=="\n":
			array.append(" ")
		else:
			array.append(p.text)
	array_string = " ".join(array)
	return array_string

doc_temp = []
doc_counter = []
for items in words:
	doc_counter.append(0)
doc_substring_counter = 0
doc_number_of_docs = 0
doc_file_names_to_read = []

for doc in doc_files:
	print(doc)
	doc_counter[:] = []
	for items in words:
		doc_counter.append(0)
	doc_file_name = "C:\Python27\\files_to_parse\\"+doc
	text = getDocxText(doc_file_name) #get string of text content of doc
	text_string = text.encode("utf-8")
	
	text_file.write("\n"+doc+"\n")
	text_file.write("----------------------\n")

	doc_substring_counter = 0

	for substring in words:
		if re.search(substring,text_string,re.IGNORECASE):
			doc_counter[doc_substring_counter]+=1
		doc_substring_counter +=1	

	doc_index = 0

	for substring in words:
		text_file.write(substring + " appears " + str(doc_counter[doc_index]) + " times\n")
		doc_index+=1

	doc_value = all_same(doc_counter)	
	if doc_value == False:
		text_file.write("**Read this document**\n")
		doc_number_of_docs+=1
		doc_file_names_to_read.append(doc)

#PPT/PPTX Parser#

for eachfile in ppt_files:
	print(eachfile)
	counter[:] = []
	for items in words:
		counter.append(0)
	f = open("C:\Python27\\files_to_parse\\"+eachfile, "rb")
	prs = Presentation(f)
	text_file.write("\n"+eachfile+"\n")
	text_file.write("----------------------\n")
	for slide in prs.slides:
		for shape in slide.shapes:
			if not shape.has_text_frame:
				continue
			for paragraph in shape.text_frame.paragraphs:
				for run in paragraph.runs:
					substring_counter = 0
					for substring in words:
						if re.search(substring,run.text, re.IGNORECASE):
							counter[substring_counter]+=1
						substring_counter +=1
	index = 0

	for substring in words:
		text_file.write(substring + " appears " + str(counter[index]) + " times\n")
		index+=1

	value = all_same(counter)
	if value == False:
		text_file.write("**Read this document**\n")
		number_of_docs+=1
		file_names_to_read.append(eachfile)

#PDF Parser#

def convert(fname, pages=None):
    if not pages:
        pagenums = set()
    else:
        pagenums = set(pages)

    output = StringIO()
    manager = PDFResourceManager()
    converter = TextConverter(manager, output, laparams=LAParams())
    interpreter = PDFPageInterpreter(manager, converter)

    infile = file(fname, 'rb')
    for page in PDFPage.get_pages(infile, pagenums):
        interpreter.process_page(page)
    infile.close()
    converter.close()
    text = output.getvalue()
    output.close
    return text  

pdf_temp = []
pdf_counter = []
for items in words:
		pdf_counter.append(0)
pdf_substring_counter = 0
pdf_number_of_docs = 0
pdf_file_names_to_read = []
files_string = []

for pdf in pdf_files:
	print(pdf)
	pdf_counter[:] = []
	for items in words:
		pdf_counter.append(0)
	text_file.write("\n"+pdf+"\n")
	text_file.write("----------------------\n")
	pdf_file_name = "C:\Python27\\files_to_parse\\"+pdf
	text = convert(pdf_file_name) #get string of text content of pdf
	files_string [:] = []
	for items in text:
		if items=="\n":
			files_string.append(" ")
		else:
			files_string.append(items)

	myString = "".join(files_string)

	pdf_substring_counter=0	

	for substring in words:
		if re.search(substring,myString,re.IGNORECASE):
			pdf_counter[pdf_substring_counter]+=1
		pdf_substring_counter +=1	

	pdf_index = 0

	for substring in words:
		text_file.write(substring + " appears " + str(pdf_counter[pdf_index]) + " times\n")
		pdf_index+=1

	pdf_value = all_same(pdf_counter)
	if pdf_value == False:
		text_file.write("**Read this document**\n")
		pdf_number_of_docs+=1
		pdf_file_names_to_read.append(pdf)

#Summary#

text_file.write("\nSummary\n----------------------\nRead the following "+str(number_of_docs)+"/"+str(len(ppt_files))+ " scanned PPT/PPTX documents:\n")
for items in file_names_to_read:
	text_file.write(items)
	text_file.write("\n")

text_file.write("\nRead the following "+str(pdf_number_of_docs)+"/"+str(len(pdf_files))+ " scanned PDF documents:\n")
for items in pdf_file_names_to_read:
	text_file.write(items)
	text_file.write("\n")

text_file.write("\nRead the following "+str(doc_number_of_docs)+"/"+str(len(doc_files))+ " scanned DOCX documents:\n")
for items in doc_file_names_to_read:
	text_file.write(items)
	text_file.write("\n")

text_file.write("----------------------\n")

text_file.close()